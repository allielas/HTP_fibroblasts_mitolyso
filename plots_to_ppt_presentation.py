from pathlib import Path
import argparse

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tif", ".tiff", ".webp"}


def iter_image_files(input_dir: Path):
    """Loop through the input directory and yield image files with supported extensions.

    Args:
        input_dir (Path): The directory to search for image files.

    Yields:
        Path: The paths to the image files.
    """
    for path in sorted(input_dir.iterdir()):
        if path.is_file() and path.suffix.lower() in IMAGE_EXTENSIONS:
            yield path


def image_title_from_path(image_path: Path, max_tokens=2, use_suffix=True):
    """Extract a title from the image file name.

    Args:
        image_path (Path): The path to the image file.
        max_tokens (int, optional): The maximum number of tokens to use for the section key. Defaults to 2.
        use_suffix (bool, optional): Whether to use the suffix when extracting the title. Defaults to True.

    Returns:
        str: The title extracted from the image file name.
    """
    raw_tokens = (
        image_path.stem.split("_")
        if "_" in image_path.stem
        else image_path.stem.split(" ")
    )
    if len(raw_tokens) > 2 * max_tokens:
        max_tokens_prefix = round(len(raw_tokens) / 2) + 1
        max_tokens_suffix = round(len(raw_tokens) / 2)
    else:
        max_tokens_prefix, max_tokens_suffix = max_tokens, max_tokens

    text_section = get_section_key(image_path, max_tokens=max_tokens_prefix, first=True)
    text_suffix = get_section_key(image_path, max_tokens=max_tokens_suffix, first=False)
    if text_section is None and text_suffix is None:
        print(
            f"Warning: Could not extract section or suffix from image name {image_path.name}"
        )
        text = image_path.stem
        text = text.replace("_boxplot", "")
    else:
        text = ""
        text_section_tokens = text_section.split("_") if text_section else []
        text_suffix_tokens = text_suffix.split("_") if text_suffix else []
        all_tokens = text_section_tokens + text_suffix_tokens
        if all_tokens:
            all_tokens_unique = list(dict.fromkeys(all_tokens))
            intersecting_tokens = set(text_section_tokens).intersection(
                set(text_suffix_tokens)
            )
            if not use_suffix:
                # when not using suffix, don't add the suffix tokens to the text, but still use them to determine if a newline is needed
                all_tokens_unique = text_section_tokens
            for i, token in enumerate(all_tokens_unique):
                # add a newline if the previous token overlaps with the previous section or suffix, otherwise add a space
                if i == 0:
                    text += token
                elif token in intersecting_tokens:
                    text += "\n" + token
                    # for making titles that don't use the suffix tokens but we want the newline for formatting
                else:
                    text += " " + token
        else:
            text = (
                f"{text_section}\n{text_suffix}"
                if text_section and text_suffix
                else None
            )
    return text


def add_image_to_slide(slide, image_path: Path, prs: Presentation, max_tokens=2):
    """Add one image to a blank slide and scale it to fit while preserving aspect ratio.

    Args:
        slide (Presentation.Slide): The slide to add the image to.
        image_path (Path): The path to the image file.
        prs (Presentation): The presentation object.
        max_tokens (int, optional): The maximum number of tokens to use for the section key. Defaults to 2.

    Returns:
        None
    """
    # calculate the available width for the image after accounting for the label box and margins
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    label_width = int(slide_width * 0.22)
    image_left_margin = Inches(0.25)
    image_right_margin = Inches(0.25)
    image_left = label_width + image_left_margin
    image_width_available = slide_width - image_left - image_right_margin
    # add the label box to the left of the image and set its properties
    label_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.6), label_width, slide_height - Inches(1.2)
    )
    label_frame = label_box.text_frame
    label_frame.clear()
    label_frame.word_wrap = True
    label_frame.auto_size = None
    label_paragraph = label_frame.paragraphs[0]
    label_paragraph.alignment = PP_ALIGN.LEFT
    label_run = label_paragraph.add_run()

    text = image_title_from_path(image_path, max_tokens=max_tokens)
    # add the text
    label_run.text = text

    # add the image to the slide and scale it to fit while preserving aspect ratio
    picture = slide.shapes.add_picture(str(image_path), 0, 0)

    image_ratio = picture.width / picture.height
    image_ratio_available = image_width_available / slide_height

    if image_ratio > image_ratio_available:
        picture.width = image_width_available
        picture.height = int(image_width_available / image_ratio)
        picture.left = image_left
        picture.top = int((slide_height - picture.height) / 2)
    else:
        picture.height = slide_height
        picture.width = int(slide_height * image_ratio)
        picture.top = 0
        picture.left = int(image_left + (image_width_available - picture.width) / 2)


def get_section_key(image_path: Path, max_tokens=2, first=True):
    """Return the first two underscore-separated tokens from the image name, if present

    Args:
        image_path (Path): The path to the image file.
        max_tokens (int, optional): The maximum number of tokens to return. Defaults to 2.

    Returns:
        str: The section key, or None if not found.
    """
    tokens = image_path.stem.split("_")
    if len(tokens) < max_tokens:
        return None
    if "boxplot" in tokens:
        tokens.remove("boxplot")
    if first:
        return "_".join(tokens[:max_tokens])
    else:
        return "_".join(tokens[-max_tokens:])


def add_section_divider_slide(prs: Presentation, section_title: str):
    """Add a blank divider slide with large section text.

    Args:
        prs (Presentation): The presentation object.
        section_title (str): The title for the section divider slide.
    """
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    left = Inches(0.5)
    top = Inches(2.5)
    width = prs.slide_width - Inches(1.0)
    height = Inches(1.5)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.clear()
    text_frame.word_wrap = True

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = section_title
    run.font.size = Pt(46)
    run.font.bold = True


def images_to_pptx(input_dir: Path, output_pptx: Path, ntokens=2):
    """Create a PowerPoint presentation from image files in a directory.

    Args:
        input_dir (Path): The directory containing image files.
        output_pptx (Path): The path to the output PowerPoint file.
        ntokens (int, optional): The number of tokens to use for section keys. Defaults to 2.

    Raises:
        ValueError: If no supported image files are found in the input directory.
    """
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    image_files = list(iter_image_files(input_dir))
    if not image_files:
        raise ValueError(f"No supported image files found in {input_dir}")

    previous_section_key = None
    # add a section heading slide if the first X tokens of the image name are different from the previous image
    for image_path in image_files:
        # also use more tokens if the title is long to make more section slides for better organization
        raw_tokens = (
            image_path.stem.split("_")
            if "_" in image_path.stem
            else image_path.stem.split(" ")
        )
        if len(raw_tokens) > 2 * ntokens:
            ntokens_prefix = round(len(raw_tokens) / 2) + 1
        else:
            ntokens_prefix = ntokens
        print(
            f"image_path: {image_path}, len raw_tokens: {len(raw_tokens)}, ntokens_prefix: {ntokens_prefix}"
        )
        current_section_key = get_section_key(
            image_path, max_tokens=ntokens_prefix, first=True
        )
        if (
            current_section_key is not None
            and previous_section_key is not None
            and current_section_key != previous_section_key
        ):
            current_section_text = image_title_from_path(
                image_path, max_tokens=ntokens, use_suffix=False
            )
            add_section_divider_slide(prs, str(current_section_text))

        # add slide to ppt
        slide = prs.slides.add_slide(blank_slide_layout)
        add_image_to_slide(slide, image_path, prs, max_tokens=ntokens)
        previous_section_key = current_section_key

    prs.save(str(output_pptx))


def parse_args():
    """Parse command-line arguments.

    Returns:
        ArgumentParser: The parsed command-line arguments.
    """
    parser = argparse.ArgumentParser(
        description="Create a PowerPoint where each slide contains one image from a directory."
    )
    parser.add_argument(
        "-i",
        "--input-dir",
        type=Path,
        help="Directory containing images to add to the presentation",
    )
    parser.add_argument(
        "-o",
        "--output",
        type=Path,
        default=Path("test.pptx"),
        help="Output .pptx file name (default: test.pptx)",
    )
    parser.add_argument(
        "-n",
        "--ntokens",
        type=int,
        default=2,
        help="Number of tokens to use for section keys (default: 2)",
    )
    return parser.parse_args()


def main():
    args = parse_args()
    images_to_pptx(args.input_dir, args.output, args.ntokens)
    print(f"Saved {args.output} with images from {args.input_dir}")


if __name__ == "__main__":
    main()

# python plots_to_ppt_presentation.py -i segmentation_test/plots -o segmentation_test/test.pptx
